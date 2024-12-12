import transformers
from diffusers import StableDiffusionPipeline
from pptx import Presentation
from pptx.util import Inches

def generate_content(prompt):
    # Initialize and use an open-source language model
    model = transformers.pipeline("text-generation", model="EleutherAI/gpt-neo-1.3B")
    return model(prompt, max_length=200)[0]['generated_text']

def generate_image(prompt):
    # Initialize and use Stable Diffusion
    model = StableDiffusionPipeline.from_pretrained("CompVis/stable-diffusion-v1-4")
    image = model(prompt).images[0]
    image.save("slide_image.png")
    return "slide_image.png"

def create_presentation(title, content):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title

    # Content slides
    for topic in content:
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = topic['title']
        
        # Add text
        body_shape = slide.shapes.placeholders[1]
        body_shape.text = topic['content']
        
        # Add image
        img_path = generate_image(topic['title'])
        slide.shapes.add_picture(img_path, Inches(1), Inches(1), height=Inches(3))

    prs.save('ai_generated_presentation.pptx')

# Example usage
title = "AI in Modern Technology"
content = [
    {
        "title": "Introduction to AI",
        "content": generate_content("Write an introduction to Artificial Intelligence")
    },
    {
        "title": "Machine Learning",
        "content": generate_content("Explain Machine Learning in simple terms")
    },
    {
        "title": "Future of AI",
        "content": generate_content("Discuss the potential future developments in AI")
    }
]

create_presentation(title, content)